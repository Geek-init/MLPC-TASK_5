"""Phase 8: build the slide deck (1 title + 6 content slides).

Produces slides.pdf (matplotlib, 16:9, no external tools needed), slides.pptx
(editable, python-pptx) and slides.md (source outline). Numbers are read from
results/ so the deck always matches the recomputed Task 5 results. No bonus.
"""
from __future__ import annotations

import json

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
import matplotlib.image as mpimg

import common as C

SLIDES = C.SOLUTION_DIR / "slides"
SLIDES.mkdir(parents=True, exist_ok=True)
ASSETS = C.SOLUTION_DIR / "report" / "report_assets"


def _j(name):
    return json.loads((C.RESULTS_DIR / name).read_text())


def _find(img):
    for d in (ASSETS, C.FIGURES_DIR):
        if (d / img).exists():
            return d / img
    return None


def build_slide_specs():
    base = _j("baseline_overall.json")
    clf = _j("best_classical_overall.json")
    post = _j("best_postprocessed_overall.json")
    audit = _j("dataset_audit.json")
    gain = clf["nonhidden_test_macro_f1_official"] - base["nonhidden_test_macro_f1_official"]
    ppdelta = post["nonhidden_test_macro_f1_official"] - post["nonhidden_test_macro_f1_no_postproc"]

    return base, clf, post, audit, [
        {"kind": "title",
         "title": "MLPC 2026 Task 5\nSound Event Detection Challenge",
         "subtitle": "Group: Quantized Transformers\nAleksandr Masiev  •  Aleksandar Cvetkovic",
         "foot": "Non-bonus solution — precomputed audio features only"},
        {"kind": "bullets",
         "title": "1. Task & Final System Architecture",
         "bullets": [
             f"Goal: detect 15 domestic sound events with onsets/offsets in recordings of any length",
             f"Dataset: {audit['n_train']} train / {audit['n_validation']} val / {audit['n_test_hidden']} hidden test; "
             f"{audit['feature_dim']} precomputed features per 1s segment (hop {C.HOP_SIZE}s)",
             "Pipeline: 1s segments → standardize → multi-label classifier → per-second activity",
             "→ median-filter smoothing → merge consecutive seconds into onset/offset intervals",
             "Metric: official segment-based Macro F1 (1s resolution), evaluated with evaluate.py",
             "Honest protocol: tune on dev-validation, report on a held-out non-hidden test split"],
         "img": None},
        {"kind": "bullets",
         "title": "2. Baseline Reproduction",
         "bullets": [
             "Per-class DecisionTree (max_depth=20, max_features='sqrt'), MultiOutputClassifier",
             "Trained on 50k subsampled raw segments — reproduced exactly (seed 42)",
             f"Non-hidden test Macro F1 = {base['nonhidden_test_macro_f1_official']:.3f}",
             "Best on loud/sustained classes (running_water, vacuum_cleaner, phone_ringing)",
             "Weak on short/rare events (window_open_close, wardrobe_drawer, light_switch)",
             "Limits: no class co-occurrence, 1s resolution, under-fit on 960-dim imbalanced data"],
         "img": _find("fig_label_distribution.png")},
        {"kind": "bullets",
         "title": "3. Classical Classifiers & Hyperparameters",
         "bullets": [
             "Start from Task 4 best: one-vs-rest linear ridge (closed-form, balanced weights)",
             "Tuned 2+ hyperparameters: ridge alpha {0.1..10} × sigmoid threshold {0..0.75}",
             "Also: logistic regression (C, class_weight) and RandomForest (depth)",
             f"Best classical: {clf['model']} ({clf['params']}), threshold {clf['threshold']}",
             f"Non-hidden test Macro F1 = {clf['nonhidden_test_macro_f1_official']:.3f} "
             f"(baseline {base['nonhidden_test_macro_f1_official']:.3f}, {gain:+.3f})"],
         "img": _find("fig_hyperparam.png")},
        {"kind": "bullets",
         "title": "4. Post-Processing (Median Filter)",
         "bullets": [
             "Per-class temporal median filter on per-second predictions (window 1 = none)",
             "Removes isolated false-positive blips; fills single-second gaps",
             "Window selected on dev-validation; evaluated once on non-hidden test",
             f"Selected window = {post['selected_window']}",
             f"Non-hidden test: {post['nonhidden_test_macro_f1_no_postproc']:.3f} → "
             f"{post['nonhidden_test_macro_f1_official']:.3f} ({ppdelta:+.3f})"],
         "img": _find("fig_postproc.png")},
        {"kind": "bullets",
         "title": "5. Qualitative Error Analysis",
         "bullets": [
             "Success: dominant sustained events detected with good temporal overlap",
             "False positives: extra short transients from acoustically similar onsets",
             "Misses: short/rare events lost at 1s resolution; boundary timing shifts",
             "Visuals use the precomputed mel feature representation (no raw waveform)"],
         "img": _find("error_case_1.png")},
        {"kind": "bullets",
         "title": "6. Final System, Limits & Deployment",
         "bullets": [
             f"Final system: {clf['model']} + median filter (w={post['selected_window']}); "
             f"retrained on train + full validation for the hidden test CSV",
             "Cheap & low-latency (one matrix multiply per second) — good for on-device smart-home use",
             "Limits: independent per-class linear models, 1s resolution, label noise, imbalance",
             "Future: temporal CNN/CRNN, per-class threshold calibration, hysteresis smoothing",
             "Deployment: precision matters (false alarms erode trust); validate across rooms/devices"],
         "img": _find("fig_model_comparison.png")},
    ]


# --------------------------------------------------------------------------- #
# matplotlib PDF
# --------------------------------------------------------------------------- #
def render_pdf(specs):
    with PdfPages(SLIDES / "slides.pdf") as pdf:
        for s in specs:
            fig = plt.figure(figsize=(13.33, 7.5))
            fig.patch.set_facecolor("white")
            if s["kind"] == "title":
                fig.text(0.5, 0.62, s["title"], ha="center", va="center",
                         fontsize=30, fontweight="bold")
                fig.text(0.5, 0.40, s["subtitle"], ha="center", va="center", fontsize=17)
                fig.text(0.5, 0.12, s["foot"], ha="center", va="center",
                         fontsize=12, style="italic", color="#666")
            else:
                fig.text(0.06, 0.90, s["title"], ha="left", va="center",
                         fontsize=22, fontweight="bold", color="#283593")
                fig.add_artist(plt.Line2D([0.06, 0.94], [0.85, 0.85], color="#283593", lw=2))
                has_img = s.get("img") is not None
                text_w = 0.52 if has_img else 0.9
                y = 0.76
                for b in s["bullets"]:
                    fig.text(0.07, y, "•", fontsize=13, color="#283593")
                    fig.text(0.10, y, b, ha="left", va="top", fontsize=12.5, wrap=True)
                    y -= 0.115
                if has_img:
                    ax = fig.add_axes([0.60, 0.10, 0.36, 0.66])
                    ax.imshow(mpimg.imread(str(s["img"])))
                    ax.axis("off")
            pdf.savefig(fig)
            plt.close(fig)


# --------------------------------------------------------------------------- #
# python-pptx (editable)
# --------------------------------------------------------------------------- #
def render_pptx(specs):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:  # noqa: BLE001
        print(f"python-pptx unavailable ({exc}); skipping .pptx")
        return
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for s in specs:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
        tf = tb.text_frame
        tf.text = s["title"].replace("\n", "  ")
        tf.paragraphs[0].font.size = Pt(28 if s["kind"] != "title" else 34)
        tf.paragraphs[0].font.bold = True
        if s["kind"] == "title":
            sb = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.3), Inches(3))
            sb.text_frame.text = s["subtitle"] + "\n\n" + s["foot"]
            continue
        has_img = s.get("img") is not None
        bw = Inches(6.6) if has_img else Inches(12.3)
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), bw, Inches(5.5))
        bf = body.text_frame
        bf.word_wrap = True
        for i, b in enumerate(s["bullets"]):
            p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(15)
        if has_img:
            slide.shapes.add_picture(str(s["img"]), Inches(7.4), Inches(1.6), width=Inches(5.5))
    prs.save(str(SLIDES / "slides.pptx"))


def write_md(specs):
    lines = []
    for s in specs:
        if s["kind"] == "title":
            lines += [f"# {s['title'].splitlines()[0]}", "", s["subtitle"].replace(chr(10), "  \n"),
                      "", f"_{s['foot']}_", "", "---", ""]
        else:
            lines += [f"## {s['title']}", ""]
            lines += [f"- {b}" for b in s["bullets"]]
            if s.get("img"):
                lines += ["", f"![]({s['img'].name})"]
            lines += ["", "---", ""]
    (SLIDES / "slides.md").write_text("\n".join(lines))


def main():
    base, clf, post, audit, specs = build_slide_specs()
    render_pdf(specs)
    render_pptx(specs)
    write_md(specs)
    print(f"Slides built: {len(specs)} slides (1 title + {len(specs)-1} content) -> "
          "slides.pdf, slides.pptx, slides.md")


if __name__ == "__main__":
    main()
